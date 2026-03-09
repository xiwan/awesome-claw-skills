#!/usr/bin/env python3
"""
PPTX Translator - Translate PowerPoint files using Amazon Translate or Bedrock LLM.

Usage:
    python translate_pptx.py input.pptx --source zh --target en
    python translate_pptx.py input.pptx --source en --target zh --engine bedrock --model anthropic.claude-3-5-sonnet-20241022-v2:0
    python translate_pptx.py input.pptx --source ja --target en --engine translate --terminology terms.csv
"""

import argparse
import json
import re
import sys
from pathlib import Path

import boto3
from botocore.exceptions import ClientError
from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID

# Language code to MSO_LANGUAGE_ID mapping
LANGUAGE_CODE_TO_LANGUAGE_ID = {
    'af': MSO_LANGUAGE_ID.AFRIKAANS,
    'am': MSO_LANGUAGE_ID.AMHARIC,
    'ar': MSO_LANGUAGE_ID.ARABIC,
    'bg': MSO_LANGUAGE_ID.BULGARIAN,
    'bn': MSO_LANGUAGE_ID.BENGALI,
    'bs': MSO_LANGUAGE_ID.BOSNIAN,
    'cs': MSO_LANGUAGE_ID.CZECH,
    'da': MSO_LANGUAGE_ID.DANISH,
    'de': MSO_LANGUAGE_ID.GERMAN,
    'el': MSO_LANGUAGE_ID.GREEK,
    'en': MSO_LANGUAGE_ID.ENGLISH_US,
    'es': MSO_LANGUAGE_ID.SPANISH,
    'et': MSO_LANGUAGE_ID.ESTONIAN,
    'fi': MSO_LANGUAGE_ID.FINNISH,
    'fr': MSO_LANGUAGE_ID.FRENCH,
    'fr-CA': MSO_LANGUAGE_ID.FRENCH_CANADIAN,
    'ha': MSO_LANGUAGE_ID.HAUSA,
    'he': MSO_LANGUAGE_ID.HEBREW,
    'hi': MSO_LANGUAGE_ID.HINDI,
    'hr': MSO_LANGUAGE_ID.CROATIAN,
    'hu': MSO_LANGUAGE_ID.HUNGARIAN,
    'id': MSO_LANGUAGE_ID.INDONESIAN,
    'it': MSO_LANGUAGE_ID.ITALIAN,
    'ja': MSO_LANGUAGE_ID.JAPANESE,
    'ka': MSO_LANGUAGE_ID.GEORGIAN,
    'ko': MSO_LANGUAGE_ID.KOREAN,
    'lv': MSO_LANGUAGE_ID.LATVIAN,
    'ms': MSO_LANGUAGE_ID.MALAYSIAN,
    'nl': MSO_LANGUAGE_ID.DUTCH,
    'no': MSO_LANGUAGE_ID.NORWEGIAN_BOKMOL,
    'pl': MSO_LANGUAGE_ID.POLISH,
    'ps': MSO_LANGUAGE_ID.PASHTO,
    'pt': MSO_LANGUAGE_ID.BRAZILIAN_PORTUGUESE,
    'ro': MSO_LANGUAGE_ID.ROMANIAN,
    'ru': MSO_LANGUAGE_ID.RUSSIAN,
    'sk': MSO_LANGUAGE_ID.SLOVAK,
    'sl': MSO_LANGUAGE_ID.SLOVENIAN,
    'so': MSO_LANGUAGE_ID.SOMALI,
    'sq': MSO_LANGUAGE_ID.ALBANIAN,
    'sr': MSO_LANGUAGE_ID.SERBIAN_LATIN,
    'sv': MSO_LANGUAGE_ID.SWEDISH,
    'sw': MSO_LANGUAGE_ID.SWAHILI,
    'ta': MSO_LANGUAGE_ID.TAMIL,
    'th': MSO_LANGUAGE_ID.THAI,
    'tr': MSO_LANGUAGE_ID.TURKISH,
    'uk': MSO_LANGUAGE_ID.UKRAINIAN,
    'ur': MSO_LANGUAGE_ID.URDU,
    'vi': MSO_LANGUAGE_ID.VIETNAMESE,
    'zh': MSO_LANGUAGE_ID.CHINESE_SINGAPORE,
    'zh-TW': MSO_LANGUAGE_ID.CHINESE_HONG_KONG_SAR,
}

LANGUAGE_NAMES = {
    'en': 'English', 'zh': 'Chinese (Simplified)', 'zh-TW': 'Chinese (Traditional)',
    'ja': 'Japanese', 'ko': 'Korean', 'de': 'German', 'fr': 'French',
    'es': 'Spanish', 'pt': 'Portuguese', 'it': 'Italian', 'ru': 'Russian',
    'ar': 'Arabic', 'hi': 'Hindi', 'th': 'Thai', 'vi': 'Vietnamese',
}


class TranslationEngine:
    """Base class for translation engines."""
    
    def translate(self, text: str, source_lang: str, target_lang: str) -> str:
        raise NotImplementedError


class AmazonTranslateEngine(TranslationEngine):
    """Amazon Translate engine for traditional machine translation."""
    
    def __init__(self, terminology_names: list = None, region: str = None):
        self.client = boto3.client('translate', region_name=region)
        self.terminology_names = terminology_names or []
    
    def translate(self, text: str, source_lang: str, target_lang: str) -> str:
        if not text.strip():
            return text
        try:
            response = self.client.translate_text(
                Text=text,
                SourceLanguageCode=source_lang,
                TargetLanguageCode=target_lang,
                TerminologyNames=self.terminology_names
            )
            return response.get('TranslatedText', text)
        except ClientError as e:
            if e.response['Error']['Code'] == 'ValidationException':
                print(f"  [WARN] Invalid text, skipping: {text[:50]}...")
                return text
            raise


class BedrockLLMEngine(TranslationEngine):
    """Bedrock LLM engine for context-aware translation."""
    
    def __init__(self, model_id: str = None, region: str = None, glossary: dict = None,
                 style: str = None, batch_size: int = 20):
        self.client = boto3.client('bedrock-runtime', region_name=region)
        self.model_id = model_id or 'anthropic.claude-3-5-sonnet-20241022-v2:0'
        self.glossary = glossary or {}
        self.style = style or 'professional'
        self.batch_size = batch_size
        self._cache = {}
    
    def _build_prompt(self, texts: list, source_lang: str, target_lang: str) -> str:
        source_name = LANGUAGE_NAMES.get(source_lang, source_lang)
        target_name = LANGUAGE_NAMES.get(target_lang, target_lang)
        
        glossary_section = ""
        if self.glossary:
            glossary_items = [f"  - {k} → {v}" for k, v in self.glossary.items()]
            glossary_section = f"\n\nTerminology (use these translations consistently):\n" + "\n".join(glossary_items)
        
        texts_json = json.dumps(texts, ensure_ascii=False)
        
        return f"""Translate the following texts from {source_name} to {target_name}.

Style: {self.style}{glossary_section}

Rules:
1. Preserve any formatting markers, placeholders, or special characters
2. Maintain consistent terminology across all texts
3. Keep translations natural and fluent in the target language
4. Return ONLY a JSON array of translated strings in the same order

Input texts:
{texts_json}

Output (JSON array only):"""
    
    def _call_bedrock(self, prompt: str) -> str:
        body = json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 8192,
            "messages": [{"role": "user", "content": prompt}]
        })
        
        response = self.client.invoke_model(
            modelId=self.model_id,
            body=body,
            contentType="application/json",
            accept="application/json"
        )
        
        result = json.loads(response['body'].read())
        return result['content'][0]['text']
    
    def translate_batch(self, texts: list, source_lang: str, target_lang: str) -> list:
        """Translate a batch of texts together for context consistency."""
        if not texts:
            return []
        
        # Filter out empty texts but track indices
        non_empty = [(i, t) for i, t in enumerate(texts) if t.strip()]
        if not non_empty:
            return texts
        
        indices, to_translate = zip(*non_empty)
        
        prompt = self._build_prompt(list(to_translate), source_lang, target_lang)
        response = self._call_bedrock(prompt)
        
        # Parse JSON response
        try:
            # Extract JSON array from response
            json_match = re.search(r'\[.*\]', response, re.DOTALL)
            if json_match:
                translated = json.loads(json_match.group())
            else:
                raise ValueError("No JSON array found in response")
        except (json.JSONDecodeError, ValueError) as e:
            print(f"  [WARN] Failed to parse LLM response, falling back to original: {e}")
            return texts
        
        # Reconstruct full list with translations
        result = list(texts)
        for idx, trans in zip(indices, translated):
            result[idx] = trans
        
        return result
    
    def translate(self, text: str, source_lang: str, target_lang: str) -> str:
        """Single text translation (uses batch internally with size 1)."""
        if not text.strip():
            return text
        
        cache_key = (text, source_lang, target_lang)
        if cache_key in self._cache:
            return self._cache[cache_key]
        
        result = self.translate_batch([text], source_lang, target_lang)[0]
        self._cache[cache_key] = result
        return result


def import_terminology(client, terminology_file: str, name: str = 'pptx-translator-terminology'):
    """Import terminology file to Amazon Translate."""
    print(f"Importing terminology from {terminology_file}...")
    with open(terminology_file, 'rb') as f:
        client.import_terminology(
            Name=name,
            MergeStrategy='OVERWRITE',
            TerminologyData={'File': bytearray(f.read()), 'Format': 'CSV'}
        )
    return name


def load_glossary(glossary_file: str) -> dict:
    """Load glossary file for LLM translation (JSON or simple key=value format)."""
    glossary = {}
    path = Path(glossary_file)
    
    if path.suffix == '.json':
        with open(path) as f:
            glossary = json.load(f)
    else:
        # Simple format: source=target per line
        with open(path) as f:
            for line in f:
                line = line.strip()
                if line and '=' in line:
                    src, tgt = line.split('=', 1)
                    glossary[src.strip()] = tgt.strip()
    
    return glossary


def extract_texts_from_presentation(presentation) -> list:
    """Extract all translatable texts with their locations."""
    texts = []
    
    for slide_idx, slide in enumerate(presentation.slides):
        # Shapes with text frames
        for shape in slide.shapes:
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        for para_idx, para in enumerate(cell.text_frame.paragraphs):
                            for run_idx, run in enumerate(para.runs):
                                if run.text.strip():
                                    texts.append({
                                        'text': run.text,
                                        'location': ('table', slide_idx, shape, row_idx, cell_idx, para_idx, run_idx),
                                        'run': run
                                    })
            elif shape.has_text_frame:
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(para.runs):
                        if run.text.strip():
                            texts.append({
                                'text': run.text,
                                'location': ('shape', slide_idx, shape, para_idx, run_idx),
                                'run': run
                            })
        
        # Notes
        if slide.has_notes_slide:
            for para_idx, para in enumerate(slide.notes_slide.notes_text_frame.paragraphs):
                for run_idx, run in enumerate(para.runs):
                    if run.text.strip():
                        texts.append({
                            'text': run.text,
                            'location': ('notes', slide_idx, para_idx, run_idx),
                            'run': run
                        })
    
    return texts


def translate_presentation(presentation, engine: TranslationEngine, source_lang: str, 
                          target_lang: str, batch_mode: bool = False):
    """Translate all text in a presentation."""
    
    if batch_mode and isinstance(engine, BedrockLLMEngine):
        # Batch mode: extract all texts, translate in batches, apply back
        print("Extracting texts...")
        text_items = extract_texts_from_presentation(presentation)
        
        if not text_items:
            print("No text found to translate.")
            return
        
        print(f"Found {len(text_items)} text segments")
        
        # Translate in batches
        texts = [item['text'] for item in text_items]
        batch_size = engine.batch_size
        translated_texts = []
        
        for i in range(0, len(texts), batch_size):
            batch = texts[i:i + batch_size]
            print(f"Translating batch {i//batch_size + 1}/{(len(texts) + batch_size - 1)//batch_size}...")
            translated_batch = engine.translate_batch(batch, source_lang, target_lang)
            translated_texts.extend(translated_batch)
        
        # Apply translations
        print("Applying translations...")
        for item, translated in zip(text_items, translated_texts):
            item['run'].text = translated
    else:
        # Sequential mode: translate one by one
        total_slides = len(presentation.slides)
        for slide_idx, slide in enumerate(presentation.slides, start=1):
            print(f"Slide {slide_idx}/{total_slides}")
            
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            translate_text_frame(cell.text_frame, engine, source_lang, target_lang)
                elif shape.has_text_frame:
                    translate_text_frame(shape.text_frame, engine, source_lang, target_lang)
            
            if slide.has_notes_slide:
                translate_text_frame(slide.notes_slide.notes_text_frame, engine, source_lang, target_lang)


def translate_text_frame(text_frame, engine: TranslationEngine, source_lang: str, target_lang: str):
    """Translate all text in a text frame."""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                run.text = engine.translate(run.text, source_lang, target_lang)


def main():
    parser = argparse.ArgumentParser(
        description='Translate PPTX files using Amazon Translate or Bedrock LLM'
    )
    parser.add_argument('input_file', help='Input PPTX file path')
    parser.add_argument('--source', '-s', required=True, help='Source language code (e.g., en, zh, ja)')
    parser.add_argument('--target', '-t', required=True, help='Target language code (e.g., en, zh, ja)')
    parser.add_argument('--output', '-o', help='Output file path (default: input-{target}.pptx)')
    parser.add_argument('--engine', '-e', choices=['translate', 'bedrock'], default='translate',
                       help='Translation engine: translate (Amazon Translate) or bedrock (LLM)')
    parser.add_argument('--model', '-m', help='Bedrock model ID (default: claude-3.5-sonnet)')
    parser.add_argument('--terminology', help='Terminology CSV file for Amazon Translate')
    parser.add_argument('--glossary', '-g', help='Glossary file for Bedrock LLM (JSON or key=value format)')
    parser.add_argument('--style', help='Translation style for LLM (default: professional)')
    parser.add_argument('--batch-size', type=int, default=20, help='Batch size for LLM translation')
    parser.add_argument('--region', help='AWS region')
    parser.add_argument('--no-batch', action='store_true', help='Disable batch mode for LLM')
    
    args = parser.parse_args()
    
    # Validate input file
    input_path = Path(args.input_file)
    if not input_path.exists():
        print(f"Error: Input file not found: {args.input_file}")
        sys.exit(1)
    
    # Determine output path
    if args.output:
        output_path = Path(args.output)
    else:
        output_path = input_path.with_name(f"{input_path.stem}-{args.target}.pptx")
    
    # Initialize engine
    if args.engine == 'bedrock':
        glossary = load_glossary(args.glossary) if args.glossary else {}
        engine = BedrockLLMEngine(
            model_id=args.model,
            region=args.region,
            glossary=glossary,
            style=args.style or 'professional',
            batch_size=args.batch_size
        )
        batch_mode = not args.no_batch
        print(f"Using Bedrock LLM engine: {engine.model_id}")
    else:
        terminology_names = []
        if args.terminology:
            client = boto3.client('translate', region_name=args.region)
            name = import_terminology(client, args.terminology)
            terminology_names = [name]
        engine = AmazonTranslateEngine(terminology_names=terminology_names, region=args.region)
        batch_mode = False
        print("Using Amazon Translate engine")
    
    # Load and translate presentation
    print(f"Loading {args.input_file}...")
    presentation = Presentation(args.input_file)
    
    print(f"Translating from {args.source} to {args.target}...")
    translate_presentation(presentation, engine, args.source, args.target, batch_mode=batch_mode)
    
    # Save output
    print(f"Saving {output_path}...")
    presentation.save(str(output_path))
    print("Done!")


if __name__ == '__main__':
    main()
