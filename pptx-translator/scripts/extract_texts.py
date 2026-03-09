#!/usr/bin/env python3
"""
Extract all translatable texts from a PPTX file for review or external translation.

Usage:
    python extract_texts.py input.pptx                    # Print to stdout
    python extract_texts.py input.pptx --output texts.json  # Save as JSON
    python extract_texts.py input.pptx --format csv       # Save as CSV
"""

import argparse
import csv
import json
import sys
from pathlib import Path

from pptx import Presentation


def extract_texts(pptx_path: str) -> list:
    """Extract all texts with location metadata."""
    presentation = Presentation(pptx_path)
    texts = []
    
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        # Regular shapes
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        for para_idx, para in enumerate(cell.text_frame.paragraphs):
                            full_text = ''.join(run.text for run in para.runs)
                            if full_text.strip():
                                texts.append({
                                    'id': len(texts) + 1,
                                    'slide': slide_idx,
                                    'type': 'table',
                                    'location': f"shape{shape_idx}_row{row_idx}_cell{cell_idx}_para{para_idx}",
                                    'text': full_text
                                })
            elif shape.has_text_frame:
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    full_text = ''.join(run.text for run in para.runs)
                    if full_text.strip():
                        texts.append({
                            'id': len(texts) + 1,
                            'slide': slide_idx,
                            'type': 'shape',
                            'location': f"shape{shape_idx}_para{para_idx}",
                            'text': full_text
                        })
        
        # Notes
        if slide.has_notes_slide:
            for para_idx, para in enumerate(slide.notes_slide.notes_text_frame.paragraphs):
                full_text = ''.join(run.text for run in para.runs)
                if full_text.strip():
                    texts.append({
                        'id': len(texts) + 1,
                        'slide': slide_idx,
                        'type': 'notes',
                        'location': f"notes_para{para_idx}",
                        'text': full_text
                    })
    
    return texts


def main():
    parser = argparse.ArgumentParser(description='Extract texts from PPTX file')
    parser.add_argument('input_file', help='Input PPTX file')
    parser.add_argument('--output', '-o', help='Output file path')
    parser.add_argument('--format', '-f', choices=['json', 'csv', 'text'], default='json',
                       help='Output format')
    
    args = parser.parse_args()
    
    if not Path(args.input_file).exists():
        print(f"Error: File not found: {args.input_file}", file=sys.stderr)
        sys.exit(1)
    
    texts = extract_texts(args.input_file)
    
    if args.format == 'json':
        output = json.dumps(texts, ensure_ascii=False, indent=2)
    elif args.format == 'csv':
        import io
        buffer = io.StringIO()
        writer = csv.DictWriter(buffer, fieldnames=['id', 'slide', 'type', 'location', 'text'])
        writer.writeheader()
        writer.writerows(texts)
        output = buffer.getvalue()
    else:
        output = '\n'.join(f"[{t['id']}] Slide {t['slide']} ({t['type']}): {t['text']}" for t in texts)
    
    if args.output:
        with open(args.output, 'w', encoding='utf-8') as f:
            f.write(output)
        print(f"Extracted {len(texts)} texts to {args.output}")
    else:
        print(output)


if __name__ == '__main__':
    main()
